"""Presentation check service.

Accepts a URL or a local file path to a PPTX or PDF presentation. Validates:
- the presentation is actually present and readable;
- the number of slides;
- coverage of the required structure (problem, solution, target audience,
  tech stack, demo, team, contacts) — the section list is configured per
  hackathon by the organizer.
"""
from __future__ import annotations

import io
import logging
import os
from pathlib import Path
from urllib.parse import urlparse

import httpx

logger = logging.getLogger(__name__)

# Fallback section list (matches the model default) used when a hackathon has no
# configured sections. Each entry: {"key", "label", "keywords": [...]}.
DEFAULT_SECTIONS: list[dict] = [
    {"key": "problem", "label": "Проблема", "keywords": ["проблема", "problem"]},
    {"key": "solution", "label": "Решение", "keywords": ["решение", "solution"]},
    {"key": "audience", "label": "Целевая аудитория", "keywords": ["целевая аудитория", "аудитория", "target audience", "audience"]},
    {"key": "stack", "label": "Технологический стек", "keywords": ["стек", "технологический стек", "tech stack", "технологии", "stack"]},
    {"key": "demo", "label": "Демо", "keywords": ["демо", "demo"]},
    {"key": "team", "label": "Команда", "keywords": ["команда", "team"]},
    {"key": "contacts", "label": "Контакты", "keywords": ["контакты", "contacts"]},
]


def _fetch(url: str) -> bytes | None:
    try:
        with httpx.Client(timeout=30.0, follow_redirects=True) as client:
            r = client.get(url)
            r.raise_for_status()
            return r.content
    except Exception as exc:
        logger.warning("presentation fetch failed %s: %s", url, exc)
        return None


def _load(source: str) -> tuple[bytes | None, str]:
    """Return (bytes, error). Supports local file paths and http(s) URLs."""
    if os.path.exists(source):
        try:
            return Path(source).read_bytes(), ""
        except Exception as exc:  # pragma: no cover - fs errors
            return None, f"failed to read presentation: {exc}"
    if not source.lower().startswith(("http://", "https://")):
        return None, "source is neither a local file nor an http(s) url"
    data = _fetch(source)
    if not data:
        return None, "failed to fetch presentation"
    return data, ""


def _fmt(source: str) -> str:
    p = urlparse(source).path.lower() if "://" in source else source.lower()
    if p.endswith(".pdf"):
        return "pdf"
    return "pptx"


def _slides_pdf(data: bytes) -> list[str]:
    try:
        from PyPDF2 import PdfReader  # type: ignore

        reader = PdfReader(io.BytesIO(data))
        return [(p.extract_text() or "") for p in reader.pages]
    except Exception as exc:
        logger.warning("pdf slide extract failed: %s", exc)
        return []


def _slides_pptx(data: bytes) -> list[str]:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(data))
        out: list[str] = []
        for slide in prs.slides:
            chunks: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                chunks.append(run.text)
            out.append("\n".join(chunks))
        return out
    except Exception as exc:
        logger.warning("pptx extract failed: %s", exc)
        return []


def _section_hits(slides_text: list[str], sections: list[dict]) -> list[str]:
    """Return the keys of the configured sections found in the slide text."""
    joined = "\n".join(slides_text).lower()
    found: list[str] = []
    for sec in sections:
        if not isinstance(sec, dict):
            continue
        key = sec.get("key")
        keywords = [str(k).lower() for k in (sec.get("keywords") or []) if str(k).strip()]
        if key and keywords and any(k in joined for k in keywords):
            found.append(key)
    return found


def run_presentation_check(source: str | None, sections: list[dict] | None = None) -> dict:
    """Check a presentation file/URL against a configurable section list.

    `sections` is the organizer-defined structure for the hackathon; when empty
    the built-in default (problem/solution/.../contacts) is used.
    """
    sections = sections or DEFAULT_SECTIONS
    total = len([s for s in sections if isinstance(s, dict) and s.get("key")])
    if not source:
        return {
            "status": "skipped",
            "score": 0.0,
            "slide_count": 0,
            "sections_found": [],
            "fmt": "",
            "message": "no presentation file or url provided",
        }
    data, err = _load(source)
    fmt = _fmt(source)
    if not data:
        return {
            "status": "error",
            "score": 0.0,
            "slide_count": 0,
            "sections_found": [],
            "fmt": fmt,
            "message": err or "failed to load presentation",
        }
    slides = _slides_pdf(data) if fmt == "pdf" else _slides_pptx(data)
    count = len(slides)
    found = _section_hits(slides, sections)

    # slide count: up to 3 points (ideal 8..15 slides)
    score = 0.0
    if 8 <= count <= 15:
        score += 3.0
    elif 5 <= count <= 20:
        score += 1.5
    # structure coverage: up to 7 points, proportional to required sections
    if total > 0:
        score += min(7.0, 7.0 * len(found) / total)
    score = min(10.0, round(score, 2))

    missing = total - len(found)
    if count == 0:
        msg = "не удалось распознать слайды"
    elif missing > 0:
        msg = f"{count} слайдов · разделов {len(found)}/{total}"
    else:
        msg = f"{count} слайдов · все разделы на месте"

    return {
        "status": "done",
        "score": score,
        "slide_count": count,
        "sections_found": found,
        "fmt": fmt,
        "message": msg,
    }
